#!/usr/bin/env python3
"""Generate a template-based leadership roadmap matrix from a use-case CSV.

The matrix uses:

* capability/area values as rows;
* calendar months as columns; and
* use-case names as cell contents.

Default behavior reflects the agreed 2026 roadmap view:

* August through December 2026;
* an interactive "top N per capability area" prompt (Enter shows all);
* one master slide when more than three capability areas exist;
* detail slides containing no more than three capability areas; and
* equal-size cells throughout each matrix;
* whole-cell-first fitting that keeps every use case in a cell together;
* detail text that shrinks as needed to an 8-point floor; and
* automatic continuation detail slides only when the complete cell still
  cannot fit at that 8-point floor;
* a 0.80-inch footer exclusion zone on every generated slide; and
* a separate 0.20-inch legend band immediately above the footer.

Template behavior:

* place ``roadmap_template.pptx`` in the same folder as this script;
* the script opens that file and inherits its slide master, layouts, theme,
  dimensions, fonts, and background treatment;
* template sample slides are omitted from the output by default; and
* the source template is never overwritten.

If the template uses a different name, it is accepted when it is the only
template-like PowerPoint beside the script, or it can be selected explicitly
with ``--template``.  The matrix itself uses a neutral navy, gray, white, and
black palette.

Capability counts, cumulative delivery counts, explanatory footers, source
footers, and assumption footers are intentionally omitted from the slides.
The generated legend identifies highlighted entries as ``Key use cases`` while
leaving the full footer exclusion zone untouched for template/footer content.

Example:
    python use_case_delivery_matrix_v5_cell_complete.py roadmap.csv

Non-interactive example:
    python use_case_delivery_matrix_v5_cell_complete.py roadmap.csv --top 25

Dependency:
    python -m pip install python-pptx
"""

from __future__ import annotations

import argparse
import calendar
import csv
import re
import sys
from collections import defaultdict
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path
from typing import Iterable, Sequence

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover
    print(
        "Missing dependency 'python-pptx'. Install it with:\n"
        "    python -m pip install python-pptx",
        file=sys.stderr,
    )
    raise SystemExit(3) from exc


# ---------------------------------------------------------------------------
# Visual system.  Added text inherits the template's theme font unless
# --font-family is supplied.  Geometry is scaled to the template slide size.

BASE_SLIDE_WIDTH_IN = 13.333333
BASE_SLIDE_HEIGHT_IN = 7.5
DEFAULT_FOOTER_RESERVE_IN = 0.80
DEFAULT_LEGEND_BAND_HEIGHT_IN = 0.20
DEFAULT_DETAIL_MIN_FONT_PT = 8.0

NAVY = "17324D"
BLUE = "294A66"
MID_BLUE = "3A566E"
CHARCOAL = "20262D"
SLATE = "65717C"
GRID = "CCD3D9"
HEADER_FILL = "E8EDF1"
ALT_FILL = "F4F6F7"
WHITE = "FFFFFF"
HIGHLIGHT_FILL = "DCE6EF"
HIGHLIGHT_LINE = "365A78"

AREA_PALETTE = [NAVY, BLUE, MID_BLUE, "4B6072", "596975", "687681"]
FONT_FAMILY: str | None = None

SLIDE_WIDTH_IN = BASE_SLIDE_WIDTH_IN
SLIDE_HEIGHT_IN = BASE_SLIDE_HEIGHT_IN
SCALE_X = 1.0
SCALE_Y = 1.0
LEFT = 0.18
RIGHT = 0.18
TABLE_TOP = 0.62
TABLE_BOTTOM = BASE_SLIDE_HEIGHT_IN - DEFAULT_FOOTER_RESERVE_IN - DEFAULT_LEGEND_BAND_HEIGHT_IN
TABLE_WIDTH = SLIDE_WIDTH_IN - LEFT - RIGHT
MONTH_HEADER_HEIGHT = 0.55
AREA_WIDTH = 2.30
FOOTER_RESERVE_IN = DEFAULT_FOOTER_RESERVE_IN
LEGEND_BAND_HEIGHT_IN = DEFAULT_LEGEND_BAND_HEIGHT_IN
FOOTER_TOP = BASE_SLIDE_HEIGHT_IN - FOOTER_RESERVE_IN
LEGEND_TOP = FOOTER_TOP - LEGEND_BAND_HEIGHT_IN
LEGEND_BOTTOM = FOOTER_TOP
SLIDE_LAYOUT_INDEX = 6

DEFAULT_START_MONTH = date(2026, 8, 1)
DEFAULT_END_MONTH = date(2026, 12, 1)


@dataclass(frozen=True)
class UseCase:
    name: str
    delivery_date: date
    area: str
    rank_value: tuple[int, object]
    sort_value: tuple[int, object]
    flags: tuple[str, ...]
    source_order: int


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    if len(value) != 6:
        raise ValueError(f"Expected a six-digit hex color, got {hex_value!r}")
    return RGBColor.from_string(value.upper())


def configure_geometry(
    presentation: Presentation, args: argparse.Namespace
) -> None:
    """Scale the composition while preserving exact footer/legend reserves."""
    global SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN, SCALE_X, SCALE_Y
    global LEFT, RIGHT, TABLE_TOP, TABLE_BOTTOM, TABLE_WIDTH
    global MONTH_HEADER_HEIGHT, AREA_WIDTH
    global FOOTER_RESERVE_IN, LEGEND_BAND_HEIGHT_IN
    global FOOTER_TOP, LEGEND_TOP, LEGEND_BOTTOM

    SLIDE_WIDTH_IN = float(presentation.slide_width / Inches(1))
    SLIDE_HEIGHT_IN = float(presentation.slide_height / Inches(1))
    if SLIDE_WIDTH_IN <= 0 or SLIDE_HEIGHT_IN <= 0:
        raise ValueError("The template has invalid slide dimensions")

    SCALE_X = SLIDE_WIDTH_IN / BASE_SLIDE_WIDTH_IN
    SCALE_Y = SLIDE_HEIGHT_IN / BASE_SLIDE_HEIGHT_IN
    LEFT = 0.18 * SCALE_X
    RIGHT = 0.18 * SCALE_X
    TABLE_TOP = 0.62 * SCALE_Y
    # These reserves are physical inches, not proportional scaling.  This
    # guarantees that the requested footer clearance remains exactly the same
    # on widescreen, standard, or custom-size source templates.
    FOOTER_RESERVE_IN = args.footer_reserve
    LEGEND_BAND_HEIGHT_IN = args.legend_band_height
    FOOTER_TOP = SLIDE_HEIGHT_IN - FOOTER_RESERVE_IN
    LEGEND_BOTTOM = FOOTER_TOP
    LEGEND_TOP = LEGEND_BOTTOM - LEGEND_BAND_HEIGHT_IN
    TABLE_BOTTOM = LEGEND_TOP
    TABLE_WIDTH = SLIDE_WIDTH_IN - LEFT - RIGHT
    MONTH_HEADER_HEIGHT = 0.55 * SCALE_Y
    AREA_WIDTH = min(2.30 * SCALE_X, TABLE_WIDTH * 0.32)

    if FOOTER_TOP <= 0 or LEGEND_TOP <= TABLE_TOP:
        raise ValueError(
            "The footer and legend reserves leave no usable roadmap canvas"
        )
    if TABLE_WIDTH <= 4.0 or TABLE_BOTTOM - TABLE_TOP <= 2.5:
        raise ValueError(
            "The template's slide size is too small for the roadmap matrix"
        )


def remove_all_slides(presentation: Presentation) -> None:
    """Remove template example slides while retaining masters and layouts."""
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        slide_id_list.remove(slide_id)


def remove_slide_placeholders(slide) -> None:
    """Remove cloned title/body placeholders that could overlap the matrix."""
    for placeholder in list(slide.placeholders):
        element = placeholder._element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)


def resolve_template_path(args: argparse.Namespace) -> Path:
    """Resolve an explicit template or auto-detect one beside this script."""
    script_dir = Path(__file__).resolve().parent
    output_path = Path(args.output).expanduser().resolve()

    if args.template:
        requested = Path(args.template).expanduser()
        if not requested.is_absolute():
            requested = script_dir / requested
        template_path = requested.resolve()
        if not template_path.is_file():
            raise ValueError(f"Template PowerPoint does not exist: {template_path}")
        if template_path.suffix.casefold() != ".pptx":
            raise ValueError("--template must identify a .pptx file")
        if template_path == output_path:
            raise ValueError("The output path cannot overwrite the source template")
        return template_path

    exact_default = script_dir / "roadmap_template.pptx"
    if exact_default.is_file() and exact_default.resolve() != output_path:
        return exact_default.resolve()

    candidates = sorted(
        (
            path.resolve()
            for path in script_dir.iterdir()
            if path.is_file()
            and path.suffix.casefold() == ".pptx"
            and not path.name.startswith("~$")
            and path.resolve() != output_path
        ),
        key=lambda path: path.name.casefold(),
    )
    preferred = [path for path in candidates if "template" in path.stem.casefold()]
    if len(preferred) == 1:
        return preferred[0]
    if len(preferred) > 1:
        names = ", ".join(path.name for path in preferred)
        raise ValueError(
            "Several template-named PowerPoints are beside the script. "
            f"Select one with --template: {names}"
        )

    generated_markers = (
        "delivery_matrix",
        "verified_example",
        "leadership_options",
        "_sample",
        "_test",
    )
    plausible = [
        path
        for path in candidates
        if not any(marker in path.stem.casefold() for marker in generated_markers)
    ]
    if len(plausible) == 1:
        return plausible[0]
    if len(candidates) == 1:
        return candidates[0]
    if not candidates:
        raise ValueError(
            "No template PowerPoint was found beside the script. Place "
            "roadmap_template.pptx there, or pass --template filename.pptx."
        )

    names = ", ".join(path.name for path in candidates)
    raise ValueError(
        "The template is ambiguous because several .pptx files are beside the "
        f"script. Rename the intended file roadmap_template.pptx or use --template. "
        f"Candidates: {names}"
    )


def resolve_layout_index(
    presentation: Presentation, args: argparse.Namespace
) -> int:
    """Select a layout, preferring the template's Blank layout."""
    layouts = list(presentation.slide_layouts)
    if not layouts:
        raise ValueError("The template contains no slide layouts")

    if args.layout_index is not None:
        if not 0 <= args.layout_index < len(layouts):
            raise ValueError(
                f"--layout-index must be between 0 and {len(layouts) - 1}"
            )
        return args.layout_index

    if args.layout_name:
        requested = args.layout_name.strip().casefold()
        exact = [
            index
            for index, layout in enumerate(layouts)
            if (layout.name or "").strip().casefold() == requested
        ]
        if len(exact) == 1:
            return exact[0]
        partial = [
            index
            for index, layout in enumerate(layouts)
            if requested in (layout.name or "").strip().casefold()
        ]
        if len(partial) == 1:
            return partial[0]
        available = ", ".join(
            f"{index}:{layout.name or '(unnamed)'}"
            for index, layout in enumerate(layouts)
        )
        raise ValueError(
            f"Layout {args.layout_name!r} was not uniquely found. Available: "
            + available
        )

    blank = [
        index
        for index, layout in enumerate(layouts)
        if (layout.name or "").strip().casefold() == "blank"
    ]
    if blank:
        return blank[0]

    blank_like = [
        index
        for index, layout in enumerate(layouts)
        if "blank" in (layout.name or "").strip().casefold()
    ]
    if blank_like:
        return blank_like[0]

    return min(range(len(layouts)), key=lambda index: len(layouts[index].placeholders))


def normalize_header(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "", value.casefold())


def resolve_column(
    headers: Sequence[str], requested: str, aliases: Iterable[str] = ()
) -> str:
    """Resolve a CSV column by exact, case-insensitive, then normalized match."""
    candidates = [requested, *aliases]
    for candidate in candidates:
        if candidate in headers:
            return candidate

    by_fold = {header.casefold(): header for header in headers}
    for candidate in candidates:
        if candidate.casefold() in by_fold:
            return by_fold[candidate.casefold()]

    by_normalized: dict[str, list[str]] = defaultdict(list)
    for header in headers:
        by_normalized[normalize_header(header)].append(header)
    for candidate in candidates:
        matches = by_normalized.get(normalize_header(candidate), [])
        if len(matches) == 1:
            return matches[0]

    raise ValueError(
        f"Column {requested!r} was not found. Available columns: {', '.join(headers)}"
    )


def parse_date(value: str) -> date:
    raw = value.strip()
    if not raw:
        raise ValueError("date is blank")

    try:
        return datetime.fromisoformat(raw.replace("Z", "+00:00")).date()
    except ValueError:
        pass

    for fmt in (
        "%m/%d/%Y",
        "%m/%d/%y",
        "%Y/%m/%d",
        "%d-%b-%Y",
        "%d %b %Y",
        "%b %d, %Y",
        "%B %d, %Y",
    ):
        try:
            return datetime.strptime(raw, fmt).date()
        except ValueError:
            continue
    raise ValueError(
        f"unsupported date {value!r}; use ISO YYYY-MM-DD or a common U.S. date format"
    )


def parse_month(value: str) -> date:
    try:
        parsed = datetime.strptime(value.strip(), "%Y-%m")
    except ValueError as exc:
        raise argparse.ArgumentTypeError(
            "month must use YYYY-MM, for example 2026-08"
        ) from exc
    return date(parsed.year, parsed.month, 1)


def month_floor(value: date) -> date:
    return date(value.year, value.month, 1)


def add_month(value: date) -> date:
    if value.month == 12:
        return date(value.year + 1, 1, 1)
    return date(value.year, value.month + 1, 1)


def month_sequence(start: date, end: date) -> list[date]:
    months: list[date] = []
    cursor = month_floor(start)
    end_month = month_floor(end)
    while cursor <= end_month:
        months.append(cursor)
        cursor = add_month(cursor)
    return months


def typed_sort_value(value: str) -> tuple[int, object]:
    """Numeric values sort first, then text, with blanks always last."""
    raw = value.strip()
    if not raw:
        return (2, "")
    try:
        return (0, float(raw))
    except ValueError:
        return (1, raw.casefold())


def is_truthy(value: str, truthy_values: set[str]) -> bool:
    return value.strip().casefold() in truthy_values


def load_use_cases(
    args: argparse.Namespace,
) -> tuple[list[UseCase], dict[str, object], list[str]]:
    input_path = Path(args.input).expanduser().resolve()
    if not input_path.is_file():
        raise ValueError(f"Input CSV does not exist: {input_path}")

    with input_path.open("r", encoding=args.encoding, newline="") as handle:
        reader = csv.DictReader(handle)
        if not reader.fieldnames:
            raise ValueError("The CSV has no header row")

        headers = [header.strip() for header in reader.fieldnames]
        name_col = resolve_column(
            headers, args.name_column, ("name", "use case", "use case name")
        )
        date_col = resolve_column(
            headers,
            args.date_column,
            ("deliveryDate", "delivery date", "target date", "planned delivery date"),
        )
        area_col = resolve_column(
            headers,
            args.area_column,
            ("lane", "capability", "capability area", "area"),
        )

        rank_col = ""
        if args.rank_column:
            rank_col = resolve_column(headers, args.rank_column, ("priority", "rank"))

        sort_col = ""
        if args.item_sort_column:
            sort_col = resolve_column(headers, args.item_sort_column)

        highlight_cols: list[str] = []
        if not args.no_highlight:
            if args.highlight_column:
                for requested in args.highlight_column:
                    resolved = resolve_column(headers, requested)
                    if resolved not in highlight_cols:
                        highlight_cols.append(resolved)
            else:
                try:
                    highlight_cols.append(resolve_column(headers, "dan-label"))
                except ValueError:
                    pass

        truthy_values = {
            value.strip().casefold()
            for value in args.highlight_values.split(",")
            if value.strip()
        }
        if not truthy_values:
            raise ValueError("--highlight-values must contain at least one value")

        column_map: dict[str, object] = {
            "name": name_col,
            "date": date_col,
            "area": area_col,
            "rank": rank_col,
            "sort": sort_col,
            "highlight": highlight_cols,
        }

        items: list[UseCase] = []
        warnings: list[str] = []
        for row_number, row in enumerate(reader, start=2):
            clean_row = {
                (key or "").strip(): (value or "").strip()
                for key, value in row.items()
            }
            try:
                name = clean_row.get(name_col, "").strip()
                area = clean_row.get(area_col, "").strip()
                if not name:
                    raise ValueError(f"{name_col} is blank")
                if not area:
                    raise ValueError(f"{area_col} is blank")

                delivery = parse_date(clean_row.get(date_col, ""))
                rank_value = (
                    typed_sort_value(clean_row.get(rank_col, ""))
                    if rank_col
                    else (0, row_number)
                )
                sort_value = (
                    typed_sort_value(clean_row.get(sort_col, ""))
                    if sort_col
                    else (0, row_number)
                )
                flags = tuple(
                    column
                    for column in highlight_cols
                    if is_truthy(clean_row.get(column, ""), truthy_values)
                )
                items.append(
                    UseCase(
                        name=name,
                        delivery_date=delivery,
                        area=area,
                        rank_value=rank_value,
                        sort_value=sort_value,
                        flags=flags,
                        source_order=row_number,
                    )
                )
            except ValueError as exc:
                message = f"Row {row_number}: {exc}"
                if args.skip_invalid:
                    warnings.append(message)
                else:
                    raise ValueError(message) from exc

    if not items:
        raise ValueError("No valid use cases were found in the CSV")
    return items, column_map, warnings


# ---------------------------------------------------------------------------
# Ordering and selection

def first_seen_areas(items: Sequence[UseCase]) -> list[str]:
    seen: set[str] = set()
    ordered: list[str] = []
    for item in items:
        if item.area not in seen:
            seen.add(item.area)
            ordered.append(item.area)
    return ordered


def order_areas(items: Sequence[UseCase], args: argparse.Namespace) -> list[str]:
    first_seen = first_seen_areas(items)
    counts: dict[str, int] = defaultdict(int)
    for item in items:
        counts[item.area] += 1

    if args.area_order:
        requested = [part.strip() for part in args.area_order.split(",") if part.strip()]
        actual_by_fold = {area.casefold(): area for area in first_seen}
        unknown = [area for area in requested if area.casefold() not in actual_by_fold]
        if unknown:
            raise ValueError(
                "--area-order contains values not present in the CSV: "
                + ", ".join(unknown)
            )
        explicit = [actual_by_fold[area.casefold()] for area in requested]
        return explicit + [area for area in first_seen if area not in explicit]

    if args.area_sort == "alphabetical":
        return sorted(first_seen, key=str.casefold)
    if args.area_sort == "count-desc":
        return sorted(first_seen, key=lambda area: (-counts[area], first_seen.index(area)))
    return first_seen


def prompt_for_top_per_area() -> int | None:
    while True:
        try:
            raw = input(
                "How many top-ranked use cases should be shown per capability area? "
                "[Enter = all]: "
            ).strip()
        except EOFError:
            return None
        if not raw:
            return None
        try:
            value = int(raw)
        except ValueError:
            print("Enter a whole number greater than zero, or press Enter for all.")
            continue
        if value < 1:
            print("Enter a whole number greater than zero, or press Enter for all.")
            continue
        return value


def sort_typed(
    items: Sequence[UseCase], attribute: str, direction: str
) -> list[UseCase]:
    """Sort typed values while keeping blanks last in either direction."""
    populated = [item for item in items if getattr(item, attribute)[0] != 2]
    blanks = [item for item in items if getattr(item, attribute)[0] == 2]
    populated.sort(
        key=lambda item: (getattr(item, attribute), item.source_order),
        reverse=direction == "desc",
    )
    blanks.sort(key=lambda item: item.source_order)
    return [*populated, *blanks]


def resolve_top(args: argparse.Namespace) -> int | None:
    if args.all:
        return None
    if args.top is not None:
        if args.top < 1:
            raise ValueError("--top must be a whole number greater than zero")
        return args.top
    if not args.no_prompt and sys.stdin.isatty():
        return prompt_for_top_per_area()
    return None


def select_top_per_area(
    visible_items: Sequence[UseCase],
    areas: Sequence[str],
    top_per_area: int | None,
    args: argparse.Namespace,
) -> list[UseCase]:
    if top_per_area is None:
        return list(visible_items)

    selected: list[UseCase] = []
    for area in areas:
        candidates = [item for item in visible_items if item.area == area]
        ranked = sort_typed(candidates, "rank_value", args.rank_direction)
        selected.extend(ranked[:top_per_area])
    return sorted(selected, key=lambda item: item.source_order)


def resolve_io_paths(args: argparse.Namespace, parser: argparse.ArgumentParser) -> None:
    if args.csv and args.input:
        parser.error("Provide the CSV either positionally or with --input, not both")

    raw_input = args.input or args.csv
    if not raw_input:
        candidates = sorted(Path.cwd().glob("*.csv"))
        if len(candidates) == 1:
            raw_input = str(candidates[0])
        elif sys.stdin.isatty():
            raw_input = input("CSV path: ").strip()
            if not raw_input:
                parser.error("No CSV path was supplied")
        else:
            parser.error("Supply a CSV path, for example: python script.py roadmap.csv")

    input_path = Path(raw_input).expanduser()
    args.input = str(input_path)
    if not args.output:
        args.output = str(
            input_path.with_name(f"{input_path.stem}_delivery_matrix_v5.pptx")
        )


# ---------------------------------------------------------------------------
# Text fitting. Wrapping is inserted explicitly at spaces so PowerPoint does
# not split a word merely because a cell is narrow.

def glyph_units(character: str) -> float:
    if character == " ":
        return 0.30
    if character in "ilI1|.,:;'`!":
        return 0.29
    if character in "mwMW@%&QO0":
        return 0.82
    if character.isupper():
        return 0.63
    if character.isdigit():
        return 0.56
    if character in "-/\\()[]":
        return 0.40
    return 0.53


def estimated_width_points(text: str, font_size: float) -> float:
    return sum(glyph_units(character) for character in text) * font_size


def wrap_at_words(text: str, width_in: float, font_size: float) -> list[str] | None:
    """Return explicit lines, or None when one unbroken token cannot fit."""
    words = text.split()
    if not words:
        return [""]
    limit = max(1.0, width_in * 72.0 * 0.94)
    if any(estimated_width_points(word, font_size) > limit for word in words):
        return None

    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        candidate = f"{current} {word}"
        if estimated_width_points(candidate, font_size) <= limit:
            current = candidate
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def cell_text_width(cell_width: float, highlighted: bool = False) -> float:
    # Bullet, gap, and label margins are removed from the cell width.
    return max(0.10, cell_width - (0.46 if highlighted else 0.40))


def item_height(
    item: UseCase, cell_width: float, font_size: float, *, compact: bool = False
) -> float | None:
    lines = wrap_at_words(
        item.name,
        cell_text_width(cell_width, bool(item.flags)),
        font_size,
    )
    if lines is None:
        return None
    leading = 1.10 if compact else 1.16
    gap = 0.018 if compact else 0.045
    return len(lines) * font_size * leading / 72.0 + gap


def required_cell_height(
    items: Sequence[UseCase],
    cell_width: float,
    font_size: float,
    *,
    compact: bool = False,
) -> float | None:
    if not items:
        return 0.0
    heights = [item_height(item, cell_width, font_size, compact=compact) for item in items]
    if any(height is None for height in heights):
        return None
    outer_padding = 0.05 if compact else 0.11
    return sum(float(height) for height in heights) + outer_padding


def fit_cell_font(
    items: Sequence[UseCase],
    cell_width: float,
    cell_height: float,
    maximum: float,
    minimum: float,
    *,
    compact: bool = False,
) -> tuple[float, bool]:
    if not items:
        return maximum, True
    size = maximum
    while size >= minimum - 1e-6:
        required = required_cell_height(
            items, cell_width, size, compact=compact
        )
        if required is not None and required <= cell_height + 1e-6:
            return round(size * 2.0) / 2.0, True
        size -= 0.5
    return minimum, False


def paginate_cell(
    items: Sequence[UseCase],
    cell_width: float,
    cell_height: float,
    minimum_font: float,
) -> list[list[UseCase]]:
    if not items:
        return [[]]

    pages: list[list[UseCase]] = []
    current: list[UseCase] = []
    for item in items:
        candidate = [*current, item]
        required = required_cell_height(candidate, cell_width, minimum_font)
        if required is not None and required <= cell_height + 1e-6:
            current = candidate
            continue
        if not current:
            raise ValueError(
                f"Use-case name {item.name!r} cannot fit in one detail cell at "
                f"{minimum_font:g} pt without splitting a word. Shorten that name or "
                "reduce --min-cell-font."
            )
        pages.append(current)
        current = [item]
        single_required = required_cell_height(current, cell_width, minimum_font)
        if single_required is None or single_required > cell_height + 1e-6:
            raise ValueError(
                f"Use-case name {item.name!r} cannot fit in one detail cell at "
                f"{minimum_font:g} pt without clipping. Shorten that name or reduce "
                "--min-cell-font."
            )
    if current:
        pages.append(current)
    return pages


def plan_cell_pages(
    items: Sequence[UseCase],
    cell_width: float,
    cell_height: float,
    minimum_font: float,
) -> list[list[UseCase]]:
    """Keep a complete cell together whenever it fits at the font floor.

    Continuation pagination is a last resort.  The planner tests the complete
    ordered cell at ``minimum_font`` before partitioning it, so a cell is never
    split merely to preserve a larger font.  When the complete set cannot fit,
    ``paginate_cell`` creates the fewest sequential continuation chunks that
    fit at the same floor.
    """
    if not items:
        return [[]]

    complete_height = required_cell_height(items, cell_width, minimum_font)
    if complete_height is not None and complete_height <= cell_height + 1e-6:
        return [list(items)]

    return paginate_cell(items, cell_width, cell_height, minimum_font)


def required_area_height(text: str, width_in: float, font_size: float) -> float | None:
    lines = wrap_at_words(text, width_in, font_size)
    if lines is None:
        return None
    return len(lines) * font_size * 1.12 / 72.0 + 0.10


def fit_area_font(
    text: str,
    width_in: float,
    height_in: float,
    maximum: float,
    minimum: float,
) -> tuple[float, list[str], bool]:
    size = maximum
    while size >= minimum - 1e-6:
        lines = wrap_at_words(text, width_in, size)
        required = required_area_height(text, width_in, size)
        if lines is not None and required is not None and required <= height_in + 1e-6:
            return round(size * 2.0) / 2.0, lines, True
        size -= 0.5
    return minimum, [text], False


def fit_single_line_font(
    text: str, width_in: float, maximum: float, minimum: float
) -> tuple[float, bool]:
    size = maximum
    while size >= minimum - 1e-6:
        if estimated_width_points(text, size) <= width_in * 72.0 * 0.95:
            return round(size * 2.0) / 2.0, True
        size -= 0.5
    return minimum, False


# ---------------------------------------------------------------------------
# PowerPoint primitives

def add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill: str,
    line: str = GRID,
    line_width: float = 0.65,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: float,
    color: str = CHARCOAL,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin_left: float = 0.0,
    margin_right: float = 0.0,
    margin_top: float = 0.0,
    margin_bottom: float = 0.0,
    word_wrap: bool = False,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = word_wrap
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = valign
    frame.margin_left = Inches(margin_left)
    frame.margin_right = Inches(margin_right)
    frame.margin_top = Inches(margin_top)
    frame.margin_bottom = Inches(margin_bottom)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    if FONT_FAMILY:
        run.font.name = FONT_FAMILY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_slide_header(slide, title: str, view_label: str) -> None:
    title_width = TABLE_WIDTH * 0.56
    label_width = TABLE_WIDTH - title_width - 0.18 * SCALE_X
    header_top = max(0.02, TABLE_TOP - 0.51 * SCALE_Y)
    header_height = 0.29 * SCALE_Y
    divider_top = TABLE_TOP - 0.13 * SCALE_Y

    title_size, title_fits = fit_single_line_font(
        title.upper(), title_width, 17.0, 11.0
    )
    if not title_fits:
        raise ValueError("--title is too long for the compact slide header")
    view_size, view_fits = fit_single_line_font(
        view_label.upper(), label_width, 11.0, 7.0
    )
    if not view_fits:
        raise ValueError("The generated view label is too long for the slide header")

    add_text(
        slide,
        title.upper(),
        LEFT,
        header_top,
        title_width,
        header_height,
        font_size=title_size,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        view_label.upper(),
        LEFT + TABLE_WIDTH - label_width,
        header_top,
        label_width,
        header_height,
        font_size=view_size,
        color=SLATE,
        bold=True,
        align=PP_ALIGN.RIGHT,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_rect(
        slide,
        LEFT,
        divider_top,
        TABLE_WIDTH,
        max(0.01, 0.014 * SCALE_Y),
        NAVY,
        NAVY,
        0,
    )


def add_key_use_case_legend(slide, args: argparse.Namespace) -> None:
    """Draw a compact highlight legend only inside the reserved legend band."""
    if args.no_highlight:
        return

    label = args.legend_label.strip()
    if not label:
        return

    band_height = LEGEND_BOTTOM - LEGEND_TOP
    if band_height <= 0:
        raise ValueError("The reserved legend band has no usable height")

    swatch_width = min(0.18 * SCALE_X, 0.24)
    swatch_height = min(0.12, band_height * 0.70)
    gap = min(0.08 * SCALE_X, 0.10)
    label_width = min(1.45 * SCALE_X, max(0.75, TABLE_WIDTH * 0.22))
    font_size, fits = fit_single_line_font(label, label_width, 8.0, 6.0)
    if not fits:
        raise ValueError(
            "--legend-label is too long for the reserved legend band; shorten it"
        )

    total_width = swatch_width + gap + label_width
    legend_left = LEFT + TABLE_WIDTH - total_width
    swatch_top = LEGEND_TOP + (band_height - swatch_height) / 2.0
    add_rect(
        slide,
        legend_left,
        swatch_top,
        swatch_width,
        swatch_height,
        HIGHLIGHT_FILL,
        HIGHLIGHT_LINE,
        0.65,
    )
    add_text(
        slide,
        label,
        legend_left + swatch_width + gap,
        LEGEND_TOP,
        label_width,
        band_height,
        font_size=font_size,
        color=CHARCOAL,
        bold=False,
        valign=MSO_ANCHOR.MIDDLE,
    )


def draw_matrix_frame(
    slide,
    months: Sequence[date],
    args: argparse.Namespace,
) -> float:
    month_width = (TABLE_WIDTH - AREA_WIDTH) / len(months)
    add_rect(
        slide,
        LEFT,
        TABLE_TOP,
        AREA_WIDTH,
        MONTH_HEADER_HEIGHT,
        NAVY,
        WHITE,
    )
    add_text(
        slide,
        args.area_header.upper(),
        LEFT + 0.14,
        TABLE_TOP,
        AREA_WIDTH - 0.28,
        MONTH_HEADER_HEIGHT,
        font_size=13.0,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )

    for month_index, month in enumerate(months):
        left = LEFT + AREA_WIDTH + month_index * month_width
        final_month = month_index == len(months) - 1 and args.highlight_final_month
        fill = BLUE if final_month else HEADER_FILL
        text_color = WHITE if final_month else NAVY
        add_rect(
            slide,
            left,
            TABLE_TOP,
            month_width,
            MONTH_HEADER_HEIGHT,
            fill,
            WHITE,
        )
        label = month.strftime("%b %Y").upper()
        font_size, fits = fit_single_line_font(
            label, max(0.1, month_width - 0.12), 14.0, 8.0
        )
        if not fits:
            raise ValueError(
                "Month columns are too narrow. Use a shorter horizon; month-window "
                "splitting is intentionally outside this 2026-focused version."
            )
        add_text(
            slide,
            label,
            left + 0.06,
            TABLE_TOP,
            month_width - 0.12,
            MONTH_HEADER_HEIGHT,
            font_size=font_size,
            color=text_color,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    return month_width


def add_cell_items(
    slide,
    items: Sequence[UseCase],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: float,
    bullet_color: str,
    compact: bool = False,
) -> None:
    if not items:
        return

    cursor = top + (0.025 if compact else 0.055)
    label_left = left + (0.20 if compact else 0.27)
    label_width = width - (0.25 if compact else 0.34)
    for item in items:
        lines = wrap_at_words(
            item.name,
            cell_text_width(width, bool(item.flags)),
            font_size,
        )
        if lines is None:
            raise ValueError(
                f"Internal wrap check failed for use-case name {item.name!r}"
            )
        rendered = "\n".join(lines)
        natural_height = item_height(item, width, font_size, compact=compact)
        if natural_height is None:
            raise ValueError(f"Internal height check failed for {item.name!r}")

        bullet_width = 0.11 if compact else 0.14
        bullet_size = max(2.0, min(font_size, 12.0))
        add_text(
            slide,
            "•",
            left + (0.07 if compact else 0.10),
            cursor,
            bullet_width,
            natural_height,
            font_size=bullet_size,
            color=bullet_color,
            bold=True,
            valign=MSO_ANCHOR.TOP,
        )
        label = add_text(
            slide,
            rendered,
            label_left,
            cursor,
            label_width,
            natural_height,
            font_size=font_size,
            color=CHARCOAL,
            bold=not compact,
            valign=MSO_ANCHOR.TOP,
            margin_left=0.03 if item.flags else 0.0,
            margin_right=0.03 if item.flags else 0.0,
            margin_top=0.01 if item.flags else 0.0,
            margin_bottom=0.01 if item.flags else 0.0,
            word_wrap=False,
        )
        if item.flags:
            label.fill.solid()
            label.fill.fore_color.rgb = rgb(HIGHLIGHT_FILL)
            label.line.color.rgb = rgb(HIGHLIGHT_LINE)
            label.line.width = Pt(0.65)
        cursor += natural_height

    if cursor > top + height + 0.02:
        raise ValueError("Internal layout check failed: cell content exceeded its bounds")


def draw_area_row(
    slide,
    area: str,
    area_index: int,
    row_top: float,
    row_height: float,
    months: Sequence[date],
    month_width: float,
    cells: dict[date, list[UseCase]],
    *,
    min_font: float,
    max_font: float,
    master: bool,
) -> None:
    area_color = AREA_PALETTE[area_index % len(AREA_PALETTE)]
    add_rect(slide, LEFT, row_top, AREA_WIDTH, row_height, area_color, WHITE)

    area_max = 13.0 if master else 20.0
    area_min = 3.0 if master else 11.0
    area_font, area_lines, area_fits = fit_area_font(
        area,
        AREA_WIDTH - 0.30,
        row_height - 0.12,
        area_max,
        area_min,
    )
    if not area_fits:
        raise ValueError(
            f"Capability-area label {area!r} cannot fit without clipping. "
            "Use a shorter label or fewer capability areas."
        )
    add_text(
        slide,
        "\n".join(area_lines),
        LEFT + 0.15,
        row_top + 0.06,
        AREA_WIDTH - 0.30,
        row_height - 0.12,
        font_size=area_font,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
        word_wrap=False,
    )

    for month_index, month in enumerate(months):
        left = LEFT + AREA_WIDTH + month_index * month_width
        fill = WHITE if month_index % 2 == 0 else ALT_FILL
        add_rect(slide, left, row_top, month_width, row_height, fill, GRID)
        cell_items = cells.get(month, [])
        font_size, fits = fit_cell_font(
            cell_items,
            month_width,
            row_height,
            max_font,
            min_font,
            compact=master,
        )
        if not fits:
            raise ValueError(
                f"Internal layout check failed for {area} / {month:%b %Y}."
            )
        add_cell_items(
            slide,
            cell_items,
            left,
            row_top,
            month_width,
            row_height,
            font_size=font_size,
            bullet_color=area_color,
            compact=master,
        )


# ---------------------------------------------------------------------------
# Slide planning and rendering

def new_blank_slide(presentation: Presentation):
    slide = presentation.slides.add_slide(
        presentation.slide_layouts[SLIDE_LAYOUT_INDEX]
    )
    remove_slide_placeholders(slide)
    return slide


def horizon_label(months: Sequence[date]) -> str:
    if not months:
        return ""
    if len(months) == 1:
        return months[0].strftime("%b %Y")
    if months[0].year == months[-1].year:
        return f"{months[0]:%b}–{months[-1]:%b %Y}"
    return f"{months[0]:%b %Y}–{months[-1]:%b %Y}"


def build_matrix(
    items: Sequence[UseCase], months: Sequence[date]
) -> dict[tuple[str, date], list[UseCase]]:
    month_set = set(months)
    matrix: dict[tuple[str, date], list[UseCase]] = defaultdict(list)
    for item in items:
        month = month_floor(item.delivery_date)
        if month in month_set:
            matrix[(item.area, month)].append(item)
    return matrix


def sort_matrix_cells(
    matrix: dict[tuple[str, date], list[UseCase]], args: argparse.Namespace
) -> None:
    for key, cell_items in list(matrix.items()):
        matrix[key] = sort_typed(cell_items, "sort_value", args.item_sort_direction)


def allocate_master_row_heights(
    areas: Sequence[str],
    months: Sequence[date],
    matrix: dict[tuple[str, date], list[UseCase]],
    month_width: float,
    args: argparse.Namespace,
) -> list[float]:
    """Return equal row heights and verify that the master remains complete."""
    body_height = TABLE_BOTTOM - TABLE_TOP - MONTH_HEADER_HEIGHT
    if not areas:
        return []
    row_height = body_height / len(areas)

    for area in areas:
        minimum_area = required_area_height(area, AREA_WIDTH - 0.30, 3.0)
        if minimum_area is None:
            raise ValueError(
                f"Capability-area label {area!r} has a token too long for the master view"
            )
        if minimum_area > row_height + 1e-6:
            raise ValueError(
                f"Capability-area label {area!r} cannot fit in the equal-height "
                "master row. Shorten that label or reduce the number of areas."
            )
        for month in months:
            required = required_cell_height(
                matrix[(area, month)],
                month_width,
                args.master_min_font,
                compact=True,
            )
            if required is None:
                raise ValueError(
                    f"A use-case name in {area!r} has a token too long for the "
                    "equal-size master cell."
                )
            if required > row_height + 1e-6:
                raise ValueError(
                    "The equal-size master cells cannot contain every selected "
                    f"use case in {area!r} / {month:%b %Y}, even at "
                    f"{args.master_min_font:g} pt. Use a smaller top-N value, "
                    "shorten names, or reduce --master-min-font."
                )

    return [row_height] * len(areas)


def add_master_slide(
    presentation: Presentation,
    areas: Sequence[str],
    months: Sequence[date],
    matrix: dict[tuple[str, date], list[UseCase]],
    area_indices: dict[str, int],
    args: argparse.Namespace,
) -> None:
    slide = new_blank_slide(presentation)
    add_slide_header(slide, args.title, f"Master view • {horizon_label(months)}")
    month_width = draw_matrix_frame(slide, months, args)
    row_heights = allocate_master_row_heights(
        areas, months, matrix, month_width, args
    )

    row_top = TABLE_TOP + MONTH_HEADER_HEIGHT
    for area, row_height in zip(areas, row_heights):
        cells = {month: matrix[(area, month)] for month in months}
        draw_area_row(
            slide,
            area,
            area_indices[area],
            row_top,
            row_height,
            months,
            month_width,
            cells,
            min_font=args.master_min_font,
            max_font=args.master_max_font,
            master=True,
        )
        row_top += row_height
    add_key_use_case_legend(slide, args)


def detail_page_plan(
    group: Sequence[str],
    months: Sequence[date],
    matrix: dict[tuple[str, date], list[UseCase]],
    month_width: float,
    row_height: float,
    min_font: float,
) -> tuple[dict[tuple[str, date], list[list[UseCase]]], int]:
    plan: dict[tuple[str, date], list[list[UseCase]]] = {}
    page_count = 1
    for area in group:
        for month in months:
            pages = plan_cell_pages(
                matrix[(area, month)], month_width, row_height, min_font
            )
            plan[(area, month)] = pages
            page_count = max(page_count, len(pages))
    return plan, page_count


def add_detail_group(
    presentation: Presentation,
    group: Sequence[str],
    group_number: int,
    group_count: int,
    areas: Sequence[str],
    months: Sequence[date],
    matrix: dict[tuple[str, date], list[UseCase]],
    area_indices: dict[str, int],
    args: argparse.Namespace,
) -> int:
    month_width = (TABLE_WIDTH - AREA_WIDTH) / len(months)
    body_height = TABLE_BOTTOM - TABLE_TOP - MONTH_HEADER_HEIGHT
    row_height = body_height / len(group)
    plan, page_count = detail_page_plan(
        group,
        months,
        matrix,
        month_width,
        row_height,
        args.min_cell_font,
    )

    first_area_position = areas.index(group[0]) + 1
    last_area_position = areas.index(group[-1]) + 1
    created = 0
    for page_index in range(page_count):
        slide = new_blank_slide(presentation)
        view_label = (
            f"Detail • areas {first_area_position}–{last_area_position} of {len(areas)}"
        )
        if page_count > 1:
            view_label += f" • continuation {page_index + 1}/{page_count}"
        add_slide_header(slide, args.title, view_label)
        rendered_month_width = draw_matrix_frame(slide, months, args)

        row_top = TABLE_TOP + MONTH_HEADER_HEIGHT
        for area in group:
            cells: dict[date, list[UseCase]] = {}
            for month in months:
                pages = plan[(area, month)]
                cells[month] = pages[page_index] if page_index < len(pages) else []
            draw_area_row(
                slide,
                area,
                area_indices[area],
                row_top,
                row_height,
                months,
                rendered_month_width,
                cells,
                min_font=args.min_cell_font,
                max_font=args.max_cell_font,
                master=False,
            )
            row_top += row_height
        add_key_use_case_legend(slide, args)
        created += 1
    return created


def build_presentation(
    all_items: Sequence[UseCase],
    column_map: dict[str, object],
    args: argparse.Namespace,
) -> tuple[list[str], dict[str, object]]:
    global FONT_FAMILY, SLIDE_LAYOUT_INDEX
    del column_map  # Column mapping is reported by main; it is not shown on slides.
    start = args.start_month
    end = args.end_month
    if start > end:
        raise ValueError("--start-month must be before or equal to --end-month")
    months = month_sequence(start, end)
    if len(months) > args.max_months:
        raise ValueError(
            f"The selected horizon has {len(months)} months; this version permits "
            f"at most {args.max_months}."
        )

    areas = order_areas(all_items, args)
    if not areas:
        raise ValueError("No capability areas were found")

    visible_items = [
        item for item in all_items if start <= month_floor(item.delivery_date) <= end
    ]
    if not visible_items:
        raise ValueError(
            f"No use cases fall within {start:%b %Y} through {end:%b %Y}"
        )

    top_per_area = resolve_top(args)
    selected_items = select_top_per_area(
        visible_items, areas, top_per_area, args
    )
    matrix = build_matrix(selected_items, months)
    sort_matrix_cells(matrix, args)
    area_indices = {area: index for index, area in enumerate(areas)}

    output_path = Path(args.output).expanduser().resolve()
    if output_path.suffix.casefold() != ".pptx":
        raise ValueError("--output must end in .pptx")
    template_path = resolve_template_path(args)
    presentation = Presentation(str(template_path))
    template_slide_count = len(presentation.slides)
    SLIDE_LAYOUT_INDEX = resolve_layout_index(presentation, args)
    selected_layout_name = (
        presentation.slide_layouts[SLIDE_LAYOUT_INDEX].name
        or f"Layout {SLIDE_LAYOUT_INDEX}"
    )
    if not args.keep_template_slides:
        remove_all_slides(presentation)
    retained_template_slides = len(presentation.slides)
    configure_geometry(presentation, args)
    FONT_FAMILY = args.font_family.strip() or None

    presentation.core_properties.title = args.title
    presentation.core_properties.subject = "Capability by delivery month"
    presentation.core_properties.comments = (
        "Generated from CSV using a source PowerPoint template. Ranking is applied "
        "independently within each capability area after filtering to the displayed "
        "month range. Generated roadmap content excludes the reserved footer zone."
    )

    master_created = len(areas) > args.areas_per_detail
    if master_created:
        add_master_slide(
            presentation, areas, months, matrix, area_indices, args
        )

    detail_slide_count = 0
    groups = [
        areas[index : index + args.areas_per_detail]
        for index in range(0, len(areas), args.areas_per_detail)
    ]
    for group_number, group in enumerate(groups, start=1):
        detail_slide_count += add_detail_group(
            presentation,
            group,
            group_number,
            len(groups),
            areas,
            months,
            matrix,
            area_indices,
            args,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)

    selected_by_area = {
        area: sum(1 for item in selected_items if item.area == area) for area in areas
    }
    stats: dict[str, object] = {
        "output": output_path,
        "template": template_path,
        "template_slide_count": template_slide_count,
        "template_slides_retained": retained_template_slides,
        "layout_index": SLIDE_LAYOUT_INDEX,
        "layout_name": selected_layout_name,
        "slide_size": f"{SLIDE_WIDTH_IN:.3f} x {SLIDE_HEIGHT_IN:.3f} in",
        "footer_reserve": FOOTER_RESERVE_IN,
        "legend_band_height": LEGEND_BAND_HEIGHT_IN,
        "footer_top": FOOTER_TOP,
        "legend_top": LEGEND_TOP,
        "source_count": len(all_items),
        "visible_count": len(visible_items),
        "selected_count": len(selected_items),
        "selected_by_area": selected_by_area,
        "area_count": len(areas),
        "master_created": master_created,
        "detail_slide_count": detail_slide_count,
        "generated_slide_count": (1 if master_created else 0) + detail_slide_count,
        "slide_count": len(presentation.slides),
        "top_per_area": top_per_area,
        "horizon": f"{start:%b %Y}–{end:%b %Y}",
    }
    warnings: list[str] = []
    empty_areas = [area for area, count in selected_by_area.items() if count == 0]
    if empty_areas:
        warnings.append(
            "These capability areas have no selected use cases in the displayed "
            "horizon and are shown as empty rows: " + ", ".join(empty_areas)
        )
    return warnings, stats


# ---------------------------------------------------------------------------
# Command line

def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description=(
            "Generate a template-based PowerPoint roadmap matrix with capability-"
            "area rows, delivery-month columns, a conditional master slide, and "
            "whole-cell-first detail pagination."
        ),
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )
    parser.add_argument("csv", nargs="?", help="Input CSV path")
    parser.add_argument("--input", help="Input CSV path instead of the positional path")
    parser.add_argument(
        "--output",
        help="Output .pptx path; defaults beside the CSV as <csv>_delivery_matrix_v5.pptx",
    )
    parser.add_argument(
        "--template",
        help=(
            "Template .pptx. A relative path is resolved beside this Python file. "
            "When omitted, roadmap_template.pptx or one unambiguous template-like "
            "PowerPoint beside the script is used."
        ),
    )
    parser.add_argument(
        "--keep-template-slides",
        action="store_true",
        help="Retain the template's existing slides and append generated slides",
    )
    layout = parser.add_mutually_exclusive_group()
    layout.add_argument(
        "--layout-name",
        default="",
        help="Template slide-layout name; defaults to Blank when available",
    )
    layout.add_argument(
        "--layout-index",
        type=int,
        help="Zero-based template slide-layout index",
    )
    parser.add_argument(
        "--font-family",
        default="",
        help="Override font family; empty inherits the template theme font",
    )
    parser.add_argument("--name-column", default="Name", help="Use-case name column")
    parser.add_argument(
        "--date-column", default="deliveryDate", help="Delivery-date column"
    )
    parser.add_argument(
        "--area-column", default="lane", help="Capability/area row column"
    )
    parser.add_argument(
        "--rank-column", default="priority", help="Column used for top-N ranking"
    )
    parser.add_argument(
        "--rank-direction",
        choices=("asc", "desc"),
        default="asc",
        help="Use asc when priority 1 is highest",
    )
    selection = parser.add_mutually_exclusive_group()
    selection.add_argument(
        "--top", type=int, help="Top N use cases to retain within each capability area"
    )
    selection.add_argument(
        "--all", action="store_true", help="Show all use cases without prompting"
    )
    parser.add_argument(
        "--no-prompt",
        action="store_true",
        help="When --top is omitted, show all without prompting",
    )
    parser.add_argument(
        "--item-sort-column",
        default="priority",
        help="Column used to order names inside cells; pass an empty string for CSV order",
    )
    parser.add_argument(
        "--item-sort-direction",
        choices=("asc", "desc"),
        default="asc",
        help="Ordering direction inside each cell",
    )
    parser.add_argument(
        "--area-sort",
        choices=("source", "alphabetical", "count-desc"),
        default="source",
        help="Capability row order when --area-order is absent",
    )
    parser.add_argument(
        "--area-order",
        default="",
        help="Comma-separated explicit row order; remaining areas follow in source order",
    )
    parser.add_argument(
        "--start-month",
        type=parse_month,
        default=DEFAULT_START_MONTH,
        help="First displayed month as YYYY-MM",
    )
    parser.add_argument(
        "--end-month",
        type=parse_month,
        default=DEFAULT_END_MONTH,
        help="Last displayed month as YYYY-MM",
    )
    parser.add_argument(
        "--title", default="Use Case Delivery Roadmap", help="Compact slide header"
    )
    parser.add_argument(
        "--area-header", default="Capability Area", help="Top-left matrix header"
    )
    parser.add_argument(
        "--footer-reserve",
        type=float,
        default=DEFAULT_FOOTER_RESERVE_IN,
        help=(
            "Physical inches reserved at the bottom of every generated slide; "
            "no generated roadmap or legend content enters this zone"
        ),
    )
    parser.add_argument(
        "--legend-band-height",
        type=float,
        default=DEFAULT_LEGEND_BAND_HEIGHT_IN,
        help="Physical inches reserved for the highlight legend above the footer",
    )
    parser.add_argument(
        "--legend-label",
        default="Key use cases",
        help="Text shown beside the highlight swatch; pass an empty string to omit",
    )
    parser.add_argument(
        "--highlight-column",
        action="append",
        help=(
            "Flag column used to highlight names; repeat for multiple columns. "
            "Every flag uses the same style. If omitted, dan-label is auto-detected."
        ),
    )
    parser.add_argument(
        "--highlight-values",
        default="true,1,yes,y,x",
        help="Comma-separated values treated as flagged, case-insensitively",
    )
    parser.add_argument(
        "--no-highlight",
        action="store_true",
        help="Disable all flag-driven highlighting",
    )
    parser.add_argument(
        "--areas-per-detail",
        type=int,
        default=3,
        help="Capability rows per detail slide; accepted range is 1 through 3",
    )
    parser.add_argument(
        "--min-cell-font",
        type=float,
        default=DEFAULT_DETAIL_MIN_FONT_PT,
        help=(
            "Minimum detail-slide use-case font. The complete cell is tested at "
            "this size before any continuation slide is added"
        ),
    )
    parser.add_argument(
        "--max-cell-font", type=float, default=16.0, help="Maximum detail font"
    )
    parser.add_argument(
        "--master-min-font",
        type=float,
        default=2.5,
        help="Minimum master-view font; the master is intended for zooming",
    )
    parser.add_argument(
        "--master-max-font", type=float, default=8.5, help="Maximum master-view font"
    )
    parser.add_argument(
        "--max-months",
        type=int,
        default=12,
        help="Maximum month columns; this version does not split month windows",
    )
    parser.add_argument(
        "--highlight-final-month",
        action="store_true",
        help="Apply a blue header fill to the final displayed month",
    )
    parser.add_argument(
        "--skip-invalid", action="store_true", help="Skip invalid CSV rows"
    )
    parser.add_argument("--encoding", default="utf-8-sig", help="CSV encoding")
    return parser


def main(argv: Sequence[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)
    resolve_io_paths(args, parser)

    if not 1 <= args.areas_per_detail <= 3:
        parser.error("--areas-per-detail must be between 1 and 3")
    if min(
        args.min_cell_font,
        args.max_cell_font,
        args.master_min_font,
        args.master_max_font,
    ) <= 0:
        parser.error("font sizes must be greater than zero")
    if args.min_cell_font > args.max_cell_font:
        parser.error("--min-cell-font cannot exceed --max-cell-font")
    if args.min_cell_font < DEFAULT_DETAIL_MIN_FONT_PT:
        parser.error(
            f"--min-cell-font cannot be less than "
            f"{DEFAULT_DETAIL_MIN_FONT_PT:g} pt"
        )
    if args.master_min_font > args.master_max_font:
        parser.error("--master-min-font cannot exceed --master-max-font")
    if args.max_months < 1:
        parser.error("--max-months must be at least 1")
    if args.footer_reserve < 0:
        parser.error("--footer-reserve cannot be negative")
    if args.legend_band_height <= 0:
        parser.error("--legend-band-height must be greater than zero")

    try:
        items, column_map, input_warnings = load_use_cases(args)
        layout_warnings, stats = build_presentation(items, column_map, args)
    except (OSError, ValueError) as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2

    for warning in [*input_warnings, *layout_warnings]:
        print(f"WARNING: {warning}", file=sys.stderr)

    print(f"Created {stats['output']}")
    print(
        f"Template: {stats['template']} | "
        f"layout {stats['layout_index']} ({stats['layout_name']}) | "
        f"slide size {stats['slide_size']}"
    )
    print(
        f"Reserved zones: {stats['footer_reserve']:.2f} in footer + "
        f"{stats['legend_band_height']:.2f} in legend band | "
        f"matrix ends at y={stats['legend_top']:.2f} in | "
        f"footer begins at y={stats['footer_top']:.2f} in"
    )
    if stats["template_slides_retained"]:
        print(
            f"Template slides retained: {stats['template_slides_retained']} "
            f"of {stats['template_slide_count']}"
        )
    else:
        print(
            f"Template sample slides omitted: {stats['template_slide_count']}"
        )
    print(
        f"Slides: {stats['slide_count']} total, "
        f"{stats['generated_slide_count']} generated "
        f"(master: {'yes' if stats['master_created'] else 'no'}, "
        f"detail: {stats['detail_slide_count']})"
    )
    print(
        f"Use cases: {stats['selected_count']} selected from "
        f"{stats['visible_count']} in {stats['horizon']} "
        f"({stats['source_count']} total CSV rows)"
    )
    if stats["top_per_area"] is None:
        print("Selection: all use cases in each capability area")
    else:
        print(f"Selection: top {stats['top_per_area']} per capability area")
    print(f"Selected by capability area: {stats['selected_by_area']}")
    print(f"Mapped columns: {column_map}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
